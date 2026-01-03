#!/usr/bin/env python3
"""Create PowerPoint presentation for Multi-Agent Research System Architecture"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Colors
LANGGRAPH_BLUE = RGBColor(41, 128, 185)
COORDINATOR_PURPLE = RGBColor(142, 68, 173)
AGENT_GREEN = RGBColor(39, 174, 96)
GRADER_ORANGE = RGBColor(230, 126, 34)
TOOLS_RED = RGBColor(192, 57, 43)
STATE_TEAL = RGBColor(22, 160, 133)
DARK_TEXT = RGBColor(50, 50, 50)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(245, 245, 245)
LIGHT_BLUE = RGBColor(214, 234, 248)
LIGHT_GREEN = RGBColor(212, 239, 223)
LIGHT_ORANGE = RGBColor(250, 229, 211)


def add_block(slide, left, top, width, height, text, color, subtitle="", font_size=16):
    """Add a block with rounded corners."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(9)
        p2.font.color.rgb = WHITE
        p2.alignment = PP_ALIGN.CENTER

    return shape


def add_outlined_block(slide, left, top, width, height, text, border_color, bg_color=WHITE, font_size=11):
    """Add an outlined block."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(2)

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = DARK_TEXT
    p.alignment = PP_ALIGN.CENTER
    return shape


def add_arrow(slide, start_x, start_y, width, height, color=DARK_TEXT, direction="right"):
    """Add an arrow shape."""
    shape_type = {
        "right": MSO_SHAPE.RIGHT_ARROW,
        "down": MSO_SHAPE.DOWN_ARROW,
        "left": MSO_SHAPE.LEFT_ARROW,
        "up": MSO_SHAPE.UP_ARROW,
    }.get(direction, MSO_SHAPE.RIGHT_ARROW)

    shape = slide.shapes.add_shape(shape_type, start_x, start_y, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_title(slide, text, color=DARK_TEXT):
    """Add slide title."""
    title = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(9.4), Inches(0.5))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = color


def create_title_slide(prs):
    """Create title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Multi-Agent Research System"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = LANGGRAPH_BLUE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(9), Inches(0.8))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Architecture & Module Overview"
    p.font.size = Pt(24)
    p.font.color.rgb = DARK_TEXT
    p.alignment = PP_ALIGN.CENTER

    # Features
    features = slide.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(9), Inches(1))
    tf = features.text_frame
    p = tf.paragraphs[0]
    p.text = "LangGraph Workflow | 7 Specialist Agents | Self-Correction Loop"
    p.font.size = Pt(16)
    p.font.color.rgb = COORDINATOR_PURPLE
    p.alignment = PP_ALIGN.CENTER


def create_overview_slide(prs):
    """Create high-level system overview slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "System Overview - Two Execution Modes")

    # Mode 1: LangGraph
    add_block(slide, Inches(0.5), Inches(0.8), Inches(4.2), Inches(0.6),
              "MODE 1: LangGraph (Default)", LANGGRAPH_BLUE, font_size=14)

    mode1_features = [
        "Self-corrective workflow",
        "Hallucination checking",
        "Answer quality grading",
        "Max 3 retry iterations"
    ]
    y = Inches(1.5)
    for feat in mode1_features:
        add_outlined_block(slide, Inches(0.6), y, Inches(4), Inches(0.35), feat, LANGGRAPH_BLUE, LIGHT_BLUE)
        y += Inches(0.4)

    # Mode 2: Coordinator
    add_block(slide, Inches(5.3), Inches(0.8), Inches(4.2), Inches(0.6),
              "MODE 2: Coordinator", COORDINATOR_PURPLE, font_size=14)

    mode2_features = [
        "Classic multi-agent",
        "Tool-based delegation",
        "Direct synthesis",
        "Conversation memory"
    ]
    y = Inches(1.5)
    for feat in mode2_features:
        add_outlined_block(slide, Inches(5.4), y, Inches(4), Inches(0.35), feat, COORDINATOR_PURPLE, RGBColor(245, 238, 248))
        y += Inches(0.4)

    # Shared: Specialist Agents
    add_block(slide, Inches(2), Inches(3.3), Inches(6), Inches(0.5),
              "7 SPECIALIST AGENTS", AGENT_GREEN, font_size=14)

    agents = ["Statistics", "Biology", "Psychology", "Philosophy", "Psychiatry", "Applications", "Product Mgr"]
    x = Inches(0.5)
    for agent in agents:
        add_outlined_block(slide, x, Inches(3.9), Inches(1.25), Inches(0.4), agent, AGENT_GREEN, LIGHT_GREEN, font_size=9)
        x += Inches(1.35)

    # External Services
    add_block(slide, Inches(2), Inches(4.5), Inches(6), Inches(0.4),
              "EXTERNAL SERVICES", TOOLS_RED, font_size=12)

    services = [("OpenAI GPT-4", Inches(2.2)), ("DuckDuckGo", Inches(4.4)), ("Tavily (opt)", Inches(6.6))]
    for svc, x in services:
        add_outlined_block(slide, x, Inches(5), Inches(1.8), Inches(0.35), svc, TOOLS_RED, LIGHT_ORANGE, font_size=10)


def create_langgraph_workflow_slide(prs):
    """Create LangGraph workflow slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "LangGraph Workflow - Self-Corrective Loop", LANGGRAPH_BLUE)

    # Column 1: Input & Routing
    add_block(slide, Inches(0.3), Inches(0.7), Inches(1.8), Inches(0.45), "Question", DARK_TEXT, font_size=12)
    add_arrow(slide, Inches(0.95), Inches(1.2), Inches(0.2), Inches(0.25), DARK_TEXT, "down")
    add_block(slide, Inches(0.3), Inches(1.5), Inches(1.8), Inches(0.55), "ROUTE", LANGGRAPH_BLUE, "question_router", font_size=12)

    # Routing paths
    add_arrow(slide, Inches(2.15), Inches(1.7), Inches(0.3), Inches(0.15), LANGGRAPH_BLUE, "right")

    # Column 2: Agent Queries
    add_block(slide, Inches(2.5), Inches(0.7), Inches(1.8), Inches(0.55), "PRIMARY\nAGENT", AGENT_GREEN, font_size=11)
    add_arrow(slide, Inches(3.15), Inches(1.3), Inches(0.2), Inches(0.25), AGENT_GREEN, "down")
    add_block(slide, Inches(2.5), Inches(1.6), Inches(1.8), Inches(0.55), "SECONDARY\nAGENTS", AGENT_GREEN, font_size=11)
    add_arrow(slide, Inches(4.35), Inches(1.8), Inches(0.3), Inches(0.15), AGENT_GREEN, "right")

    # Column 3: Web Search & Synthesis
    add_block(slide, Inches(4.7), Inches(0.7), Inches(1.8), Inches(0.55), "WEB\nSEARCH", TOOLS_RED, font_size=11)
    add_arrow(slide, Inches(5.35), Inches(1.3), Inches(0.2), Inches(0.25), TOOLS_RED, "down")
    add_block(slide, Inches(4.7), Inches(1.6), Inches(1.8), Inches(0.55), "SYNTHESIZE", COORDINATOR_PURPLE, font_size=11)
    add_arrow(slide, Inches(6.55), Inches(1.8), Inches(0.3), Inches(0.15), COORDINATOR_PURPLE, "right")

    # Column 4: Quality Checks
    add_block(slide, Inches(6.9), Inches(0.7), Inches(1.8), Inches(0.55), "HALLUCINATION\nCHECK", GRADER_ORANGE, font_size=10)
    add_arrow(slide, Inches(7.55), Inches(1.3), Inches(0.2), Inches(0.25), GRADER_ORANGE, "down")
    add_block(slide, Inches(6.9), Inches(1.6), Inches(1.8), Inches(0.55), "GRADE\nANSWER", GRADER_ORANGE, font_size=11)

    # Output or Retry
    add_arrow(slide, Inches(8.75), Inches(1.8), Inches(0.3), Inches(0.15), AGENT_GREEN, "right")
    add_block(slide, Inches(8.2), Inches(2.4), Inches(1.3), Inches(0.45), "OUTPUT", AGENT_GREEN, font_size=12)

    # Retry loop
    add_block(slide, Inches(6.9), Inches(2.4), Inches(1.2), Inches(0.45), "REFINE", GRADER_ORANGE, font_size=11)

    # Curved retry arrow (simulated with line)
    retry_label = slide.shapes.add_textbox(Inches(3.5), Inches(2.5), Inches(3), Inches(0.3))
    tf = retry_label.text_frame
    p = tf.paragraphs[0]
    p.text = "< - - - - RETRY (max 3) - - - - <"
    p.font.size = Pt(10)
    p.font.color.rgb = GRADER_ORANGE
    p.alignment = PP_ALIGN.CENTER

    # Decision diamonds description
    decisions_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(3.1), Inches(9.4), Inches(2.1)
    )
    decisions_box.fill.solid()
    decisions_box.fill.fore_color.rgb = LIGHT_GRAY
    decisions_box.line.color.rgb = DARK_TEXT
    decisions_box.line.width = Pt(1)

    # Decision points
    decisions = [
        ("1. Route Decision", "single_domain | cross_domain | web_search"),
        ("2. After Primary", "query_secondary? | web_search? | check_hallucination"),
        ("3. Hallucination", "grounded -> grade | not_grounded -> web_search"),
        ("4. Answer Grade", "useful -> output | not_useful -> refine & retry"),
    ]

    y = Inches(3.25)
    for title, desc in decisions:
        box = slide.shapes.add_textbox(Inches(0.5), y, Inches(9), Inches(0.4))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{title}:  {desc}"
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_TEXT
        y += Inches(0.45)


def create_specialist_agents_slide(prs):
    """Create specialist agents slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Specialist Agents - Domain Experts", AGENT_GREEN)

    # BaseAgent
    add_block(slide, Inches(3.5), Inches(0.7), Inches(3), Inches(0.7),
              "BaseAgent (Abstract)", DARK_TEXT, "OpenAI client, web_search, chat loop", font_size=14)

    # Inheritance arrow
    add_arrow(slide, Inches(4.75), Inches(1.45), Inches(0.2), Inches(0.3), DARK_TEXT, "down")

    # Agent grid
    agents = [
        ("Statistics Agent", "Inference, Regression\nBayesian, ML Theory", AGENT_GREEN),
        ("Biology Agent", "Genetics, Ecology\nMolecular, Evolution", RGBColor(46, 139, 87)),
        ("Psychology Agent", "Cognitive, Social\nClinical, Development", RGBColor(70, 130, 180)),
        ("Philosophy Agent", "Ethics, Epistemology\nMetaphysics, Phil. Mind", RGBColor(128, 0, 128)),
        ("Psychiatry Agent", "Disorders, Pharma\nNeuroimaging, Clinical", RGBColor(220, 20, 60)),
        ("Applications Agent", "Real-world Use\nIndustry Implementation", RGBColor(255, 140, 0)),
        ("Product Manager", "Product Strategy\nResearch-to-Product", RGBColor(0, 139, 139)),
    ]

    positions = [
        (Inches(0.3), Inches(1.9)),
        (Inches(2.55), Inches(1.9)),
        (Inches(4.8), Inches(1.9)),
        (Inches(7.05), Inches(1.9)),
        (Inches(0.3), Inches(3.2)),
        (Inches(2.55), Inches(3.2)),
        (Inches(4.8), Inches(3.2)),
    ]

    for (name, desc, color), (x, y) in zip(agents, positions):
        add_block(slide, x, y, Inches(2.15), Inches(1.1), name, color, desc, font_size=12)

    # Coordinator (separate)
    add_block(slide, Inches(7.05), Inches(3.2), Inches(2.15), Inches(1.1),
              "Coordinator", COORDINATOR_PURPLE, "Orchestrates agents\nvia delegate_to_agent", font_size=12)

    # Note
    note = slide.shapes.add_textbox(Inches(0.3), Inches(4.5), Inches(9), Inches(0.4))
    tf = note.text_frame
    p = tf.paragraphs[0]
    p.text = "Each agent has: name, description, system_prompt, web_search tool, conversation history"
    p.font.size = Pt(11)
    p.font.italic = True
    p.font.color.rgb = DARK_TEXT


def create_graders_slide(prs):
    """Create graders/quality modules slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Quality Control Modules (graders.py)", GRADER_ORANGE)

    graders = [
        ("question_router", "Classifies question to domain(s)",
         "Input: question\nOutput: primary_domain, secondary_domains, needs_web_search"),
        ("hallucination_grader", "Checks if response is grounded",
         "Input: documents, generation\nOutput: score (yes/no), unsupported_claims"),
        ("answer_grader", "Checks if answer addresses question",
         "Input: question, generation\nOutput: score (yes/no), missing_aspects"),
        ("query_refiner", "Improves query for retry",
         "Input: question, generation, issues\nOutput: refined_query, focus_areas"),
        ("response_synthesizer", "Combines multi-agent responses",
         "Input: question, agent_responses, web_results\nOutput: unified response"),
    ]

    y = Inches(0.8)
    for name, purpose, io in graders:
        # Name block
        add_block(slide, Inches(0.3), y, Inches(2.2), Inches(0.7), name, GRADER_ORANGE, font_size=11)

        # Purpose
        purpose_box = slide.shapes.add_textbox(Inches(2.7), y, Inches(3.3), Inches(0.7))
        tf = purpose_box.text_frame
        p = tf.paragraphs[0]
        p.text = purpose
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = DARK_TEXT

        # I/O
        io_box = slide.shapes.add_textbox(Inches(6), y, Inches(3.5), Inches(0.7))
        tf = io_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = io
        p.font.size = Pt(9)
        p.font.color.rgb = DARK_TEXT

        y += Inches(0.85)

    # Note about LLM
    note_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(5), Inches(9.4), Inches(0.4)
    )
    note_box.fill.solid()
    note_box.fill.fore_color.rgb = LIGHT_ORANGE
    note_box.line.fill.background()

    tf = note_box.text_frame
    p = tf.paragraphs[0]
    p.text = "All graders use GPT-3.5-turbo for speed/cost efficiency with JSON output parsing"
    p.font.size = Pt(11)
    p.font.color.rgb = DARK_TEXT
    p.alignment = PP_ALIGN.CENTER


def create_tools_slide(prs):
    """Create tools module slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Tools & Search Modules (tools.py)", TOOLS_RED)

    # Search tools
    add_block(slide, Inches(0.3), Inches(0.8), Inches(4.5), Inches(0.5),
              "SEARCH TOOLS", TOOLS_RED, font_size=14)

    tools = [
        ("web_search(query)", "DuckDuckGo search, returns Documents"),
        ("academic_search(query)", "Adds academic site filters (arxiv, pubmed, nature...)"),
        ("industry_search(query)", "Targets engineering blogs (uber, google, meta...)"),
        ("tavily_search(query)", "Premium search API (optional, needs key)"),
    ]

    y = Inches(1.4)
    for name, desc in tools:
        add_outlined_block(slide, Inches(0.4), y, Inches(2.2), Inches(0.45), name, TOOLS_RED, LIGHT_ORANGE, font_size=10)

        desc_box = slide.shapes.add_textbox(Inches(2.7), y, Inches(4.5), Inches(0.45))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = DARK_TEXT

        y += Inches(0.55)

    # Formatting tools
    add_block(slide, Inches(0.3), Inches(3.7), Inches(4.5), Inches(0.5),
              "FORMATTING HELPERS", STATE_TEAL, font_size=14)

    formatters = [
        ("format_documents_for_context(docs)", "Formats Document list for LLM"),
        ("format_agent_responses(responses)", "Formats agent dict for synthesis"),
    ]

    y = Inches(4.3)
    for name, desc in formatters:
        add_outlined_block(slide, Inches(0.4), y, Inches(3.5), Inches(0.4), name, STATE_TEAL, RGBColor(220, 240, 240), font_size=9)

        desc_box = slide.shapes.add_textbox(Inches(4), y, Inches(5.5), Inches(0.4))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = DARK_TEXT

        y += Inches(0.5)


def create_state_slide(prs):
    """Create state flow slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "ResearchState - Data Flow Through Workflow", STATE_TEAL)

    # State box
    state_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(0.7), Inches(9.4), Inches(4.5)
    )
    state_box.fill.solid()
    state_box.fill.fore_color.rgb = RGBColor(232, 245, 243)
    state_box.line.color.rgb = STATE_TEAL
    state_box.line.width = Pt(2)

    # Title in box
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(0.4))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "ResearchState (TypedDict)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = STATE_TEAL

    # State fields with arrows showing flow
    fields = [
        ("question: str", "route_question()", "User's research question"),
        ("primary_domain: str", "route_question()", "Main domain (statistics, biology...)"),
        ("secondary_domains: List[str]", "route_question()", "Cross-domain queries"),
        ("agent_responses: Dict[str, str]", "query_*_agent()", "Responses from specialists"),
        ("documents: List[Document]", "web_search_node()", "Retrieved web documents"),
        ("synthesis: str", "synthesize_responses()", "Combined response"),
        ("hallucination_grade: str", "check_hallucination()", "'grounded' or 'not_grounded'"),
        ("answer_grade: str", "grade_answer()", "'useful' or 'not_useful'"),
        ("iteration_count: int", "refine_and_retry()", "Retry counter (max 3)"),
        ("final_response: str", "generate_response()", "Final output with citations"),
    ]

    y = Inches(1.3)
    for field, setter, desc in fields:
        # Field name
        field_box = slide.shapes.add_textbox(Inches(0.5), y, Inches(3), Inches(0.35))
        tf = field_box.text_frame
        p = tf.paragraphs[0]
        p.text = field
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = DARK_TEXT

        # Setter
        setter_box = slide.shapes.add_textbox(Inches(3.5), y, Inches(2.2), Inches(0.35))
        tf = setter_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"<- {setter}"
        p.font.size = Pt(9)
        p.font.color.rgb = STATE_TEAL

        # Description
        desc_box = slide.shapes.add_textbox(Inches(5.8), y, Inches(3.5), Inches(0.35))
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(9)
        p.font.color.rgb = RGBColor(100, 100, 100)

        y += Inches(0.4)


def create_module_map_slide(prs):
    """Create module dependency map slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Module Dependency Map", DARK_TEXT)

    # Entry points
    add_block(slide, Inches(0.3), Inches(0.7), Inches(2.8), Inches(0.5),
              "ENTRY POINTS", DARK_TEXT, font_size=12)

    entries = ["main.py", "run_citation_agent.py", "run_memory_agent.py"]
    x = Inches(0.4)
    for entry in entries:
        add_outlined_block(slide, x, Inches(1.25), Inches(1.3), Inches(0.35), entry, DARK_TEXT, font_size=8)
        x += Inches(1.4)

    add_arrow(slide, Inches(2.1), Inches(1.65), Inches(0.2), Inches(0.25), DARK_TEXT, "down")

    # Core engine
    add_block(slide, Inches(0.3), Inches(2), Inches(4.2), Inches(0.45),
              "CORE ENGINE", LANGGRAPH_BLUE, font_size=12)

    core_modules = [
        ("langgraph_agent.py", "Workflow graph"),
        ("nodes.py", "Node functions"),
        ("state.py", "ResearchState"),
    ]
    x = Inches(0.4)
    for mod, desc in core_modules:
        add_outlined_block(slide, x, Inches(2.5), Inches(1.35), Inches(0.55), f"{mod}\n{desc}", LANGGRAPH_BLUE, LIGHT_BLUE, font_size=8)
        x += Inches(1.4)

    # Support modules
    add_block(slide, Inches(4.7), Inches(2), Inches(4.8), Inches(0.45),
              "SUPPORT MODULES", GRADER_ORANGE, font_size=12)

    support = [("graders.py", "Quality checks"), ("tools.py", "Search tools")]
    x = Inches(4.8)
    for mod, desc in support:
        add_outlined_block(slide, x, Inches(2.5), Inches(1.5), Inches(0.55), f"{mod}\n{desc}", GRADER_ORANGE, LIGHT_ORANGE, font_size=9)
        x += Inches(1.6)

    add_arrow(slide, Inches(4.8), Inches(3.1), Inches(0.2), Inches(0.25), DARK_TEXT, "down")

    # Agents folder
    add_block(slide, Inches(0.3), Inches(3.5), Inches(9.2), Inches(0.45),
              "agents/", AGENT_GREEN, font_size=12)

    agent_files = [
        "base_agent.py", "coordinator.py", "statistics_agent.py", "biology_agent.py",
        "psychology_agent.py", "philosophy_agent.py", "psychiatry_agent.py"
    ]
    x = Inches(0.4)
    for af in agent_files:
        add_outlined_block(slide, x, Inches(4), Inches(1.2), Inches(0.4), af, AGENT_GREEN, LIGHT_GREEN, font_size=7)
        x += Inches(1.3)

    # Advanced agents
    add_block(slide, Inches(0.3), Inches(4.55), Inches(9.2), Inches(0.45),
              "ADVANCED AGENTS", COORDINATOR_PURPLE, font_size=12)

    advanced = [
        ("citation_agent.py", "24KB"),
        ("memory_enhanced_agent.py", "10KB"),
        ("unified_research_agent.py", "20KB"),
    ]
    x = Inches(0.4)
    for af, size in advanced:
        add_outlined_block(slide, x, Inches(5.05), Inches(2.5), Inches(0.4), f"{af} ({size})", COORDINATOR_PURPLE, RGBColor(245, 238, 248), font_size=9)
        x += Inches(3)


def create_summary_slide(prs):
    """Create summary slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Architecture Summary", LANGGRAPH_BLUE)

    patterns = [
        ("State Machine", "LangGraph workflow", "Complex multi-step reasoning with conditional branching"),
        ("Self-Correction", "Hallucination check -> Retry", "Ensures quality by detecting and fixing hallucinations"),
        ("Multi-Agent Routing", "question_router", "Routes queries to appropriate domain experts"),
        ("Synthesis", "response_synthesizer", "Combines multiple expert opinions coherently"),
        ("Template Inheritance", "BaseAgent -> specialists", "Shared functionality (web search, chat loop)"),
        ("Tool Delegation", "Coordinator -> agents", "Orchestrates specialists via function calling"),
    ]

    y = Inches(0.8)
    colors = [LANGGRAPH_BLUE, GRADER_ORANGE, COORDINATOR_PURPLE, STATE_TEAL, AGENT_GREEN, TOOLS_RED]

    for (pattern, where, purpose), color in zip(patterns, colors):
        add_block(slide, Inches(0.3), y, Inches(1.8), Inches(0.55), pattern, color, font_size=11)

        where_box = slide.shapes.add_textbox(Inches(2.2), y, Inches(2.5), Inches(0.55))
        tf = where_box.text_frame
        p = tf.paragraphs[0]
        p.text = where
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = DARK_TEXT

        purpose_box = slide.shapes.add_textbox(Inches(4.7), y, Inches(5), Inches(0.55))
        tf = purpose_box.text_frame
        p = tf.paragraphs[0]
        p.text = purpose
        p.font.size = Pt(11)
        p.font.color.rgb = DARK_TEXT

        y += Inches(0.7)

    # Key insight
    insight_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(5), Inches(9.4), Inches(0.45)
    )
    insight_box.fill.solid()
    insight_box.fill.fore_color.rgb = LANGGRAPH_BLUE
    insight_box.line.fill.background()

    tf = insight_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Key: LangGraph enables self-correcting AI workflows with quality gates"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER


def main():
    """Create the full presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9

    create_title_slide(prs)
    create_overview_slide(prs)
    create_langgraph_workflow_slide(prs)
    create_specialist_agents_slide(prs)
    create_graders_slide(prs)
    create_tools_slide(prs)
    create_state_slide(prs)
    create_module_map_slide(prs)
    create_summary_slide(prs)

    output_path = "docs/Agent_Architecture.pptx"
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")


if __name__ == "__main__":
    main()
