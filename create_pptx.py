#!/usr/bin/env python3
"""Create PowerPoint presentation for Research to Application Pipeline"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Colors
RESEARCH_BLUE = RGBColor(41, 128, 185)
TRANSLATION_PURPLE = RGBColor(142, 68, 173)
PROTOTYPE_GREEN = RGBColor(39, 174, 96)
VALIDATION_ORANGE = RGBColor(230, 126, 34)
PRODUCTION_RED = RGBColor(192, 57, 43)
DARK_TEXT = RGBColor(50, 50, 50)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(245, 245, 245)

def add_phase_box(slide, left, top, width, height, text, color, subtitle=""):
    """Add a phase box with rounded corners."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(10)
        p2.font.color.rgb = WHITE
        p2.alignment = PP_ALIGN.CENTER

    return shape

def add_subitem(slide, left, top, width, height, text, border_color):
    """Add a sub-item box."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = border_color
    shape.line.width = Pt(2)

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(11)
    p.font.color.rgb = DARK_TEXT
    p.alignment = PP_ALIGN.CENTER
    return shape

def create_title_slide(prs):
    """Create title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Research to Real-World Application"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RESEARCH_BLUE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(0.8))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Translating Psychiatry Research into E-Commerce Features"
    p.font.size = Pt(24)
    p.font.color.rgb = DARK_TEXT
    p.alignment = PP_ALIGN.CENTER

    # Subtitle 2
    sub2_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(0.6))
    tf = sub2_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Optimized Checkout System"
    p.font.size = Pt(18)
    p.font.italic = True
    p.font.color.rgb = TRANSLATION_PURPLE
    p.alignment = PP_ALIGN.CENTER

    # Date
    date_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(0.5))
    tf = date_box.text_frame
    p = tf.paragraphs[0]
    p.text = "January 2026 | Multi-Agent Research System"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(128, 128, 128)
    p.alignment = PP_ALIGN.CENTER

def create_pipeline_slide(prs):
    """Create the main pipeline diagram slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9.4), Inches(0.6))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Research to Application Pipeline"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = DARK_TEXT

    # Phase boxes - horizontal layout
    box_width = Inches(1.6)
    box_height = Inches(0.7)
    start_y = Inches(1.2)

    phases = [
        ("1. RESEARCH", RESEARCH_BLUE, Inches(0.3)),
        ("2. TRANSLATE", TRANSLATION_PURPLE, Inches(2.1)),
        ("3. PROTOTYPE", PROTOTYPE_GREEN, Inches(3.9)),
        ("4. VALIDATE", VALIDATION_ORANGE, Inches(5.7)),
        ("5. PRODUCTION", PRODUCTION_RED, Inches(7.5)),
    ]

    for text, color, left in phases:
        add_phase_box(slide, left, start_y, box_width, box_height, text, color)

    # Arrows between phases
    arrow_y = Inches(1.4)
    arrow_positions = [Inches(1.95), Inches(3.75), Inches(5.55), Inches(7.35)]

    for start_x in arrow_positions:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW, start_x, arrow_y, Inches(0.15), Inches(0.15)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = DARK_TEXT
        shape.line.fill.background()

    # Sub-items for each phase
    sub_y = Inches(2.1)
    sub_height = Inches(0.4)
    sub_width = Inches(1.5)
    sub_spacing = Inches(0.45)

    # Research sub-items
    research_items = ["Psychiatry Agent", "Psychology Agent", "Statistics Agent", "Literature Review"]
    for i, item in enumerate(research_items):
        add_subitem(slide, Inches(0.35), sub_y + i * sub_spacing, sub_width, sub_height, item, RESEARCH_BLUE)

    # Translation sub-items
    trans_items = ["PM Agent", "Applications Agent", "User Personas", "Feature Definition"]
    for i, item in enumerate(trans_items):
        add_subitem(slide, Inches(2.15), sub_y + i * sub_spacing, sub_width, sub_height, item, TRANSLATION_PURPLE)

    # Prototype sub-items
    proto_items = ["UX/UI Design", "Implementation", "One-Page Checkout", "Auto-Fill System"]
    for i, item in enumerate(proto_items):
        add_subitem(slide, Inches(3.95), sub_y + i * sub_spacing, sub_width, sub_height, item, PROTOTYPE_GREEN)

    # Validation sub-items
    val_items = ["A/B Testing", "Early Stopping", "Segment Analysis", "Go/No-Go"]
    for i, item in enumerate(val_items):
        add_subitem(slide, Inches(5.75), sub_y + i * sub_spacing, sub_width, sub_height, item, VALIDATION_ORANGE)

    # Production sub-items
    prod_items = ["Phased Rollout", "Monitoring", "Iteration", "Scale"]
    for i, item in enumerate(prod_items):
        add_subitem(slide, Inches(7.55), sub_y + i * sub_spacing, sub_width, sub_height, item, PRODUCTION_RED)

    # Outcomes box at bottom
    outcomes_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(4.4), Inches(9), Inches(0.8)
    )
    outcomes_box.fill.solid()
    outcomes_box.fill.fore_color.rgb = LIGHT_GRAY
    outcomes_box.line.color.rgb = DARK_TEXT
    outcomes_box.line.width = Pt(1)

    tf = outcomes_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Target Outcomes: -20% Cart Abandonment | +10% Conversion | -67% Checkout Time | +23% CSAT"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = DARK_TEXT
    p.alignment = PP_ALIGN.CENTER

def create_research_phase_slide(prs):
    """Create detailed research phase slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.6))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Phase 1: Research Foundation"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RESEARCH_BLUE

    # Research areas
    content = [
        ("Psychiatry Research", [
            "Decision-making & Impulse Control (Bechara, 2005)",
            "Cognitive Load Theory (Sweller, 1988)",
            "Reward Systems (Berridge, 2007)"
        ]),
        ("Psychology Insights", [
            "Novelty & Habituation Effects",
            "Cognitive Biases (Anchoring, Loss Aversion)",
            "Attention & Visual Processing"
        ]),
        ("Statistical Methods", [
            "A/B Testing Framework",
            "Sample Size Calculation",
            "Early Stopping Strategies"
        ])
    ]

    y_pos = Inches(1)
    for title_text, items in content:
        section = slide.shapes.add_textbox(Inches(0.5), y_pos, Inches(9), Inches(0.4))
        tf = section.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RESEARCH_BLUE

        y_pos += Inches(0.45)

        for item in items:
            item_box = slide.shapes.add_textbox(Inches(0.7), y_pos, Inches(8.5), Inches(0.35))
            tf = item_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"• {item}"
            p.font.size = Pt(14)
            p.font.color.rgb = DARK_TEXT
            y_pos += Inches(0.35)

        y_pos += Inches(0.2)

def create_translation_phase_slide(prs):
    """Create translation phase slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.6))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Phase 2: Translation to Product"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = TRANSLATION_PURPLE

    # Left: Research Input
    left_title = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(4), Inches(0.4))
    tf = left_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Research Insights"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RESEARCH_BLUE

    left_items = [
        "Cognitive overload → abandonment",
        "Immediate gratification preference",
        "Visual attention patterns",
        "Trust signals importance"
    ]

    y = Inches(1.5)
    for item in left_items:
        box = slide.shapes.add_textbox(Inches(0.5), y, Inches(4), Inches(0.35))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = f"→ {item}"
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_TEXT
        y += Inches(0.4)

    # Right: Product Features
    right_title = slide.shapes.add_textbox(Inches(5), Inches(1), Inches(4), Inches(0.4))
    tf = right_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Product Features"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = PROTOTYPE_GREEN

    right_items = [
        "One-page checkout",
        "Express checkout options",
        "High-contrast CTAs",
        "Trust badges & SSL indicators"
    ]

    y = Inches(1.5)
    for item in right_items:
        box = slide.shapes.add_textbox(Inches(5), y, Inches(4), Inches(0.35))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = f"✓ {item}"
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_TEXT
        y += Inches(0.4)

    # Arrow in middle
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, Inches(4.2), Inches(2.2), Inches(0.6), Inches(0.4)
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = TRANSLATION_PURPLE
    arrow.line.fill.background()

    # PM Agent box
    pm_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), Inches(3.8), Inches(5), Inches(1.2)
    )
    pm_box.fill.solid()
    pm_box.fill.fore_color.rgb = TRANSLATION_PURPLE
    pm_box.line.fill.background()

    tf = pm_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Product Manager Agent"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "Bridges research insights with practical implementation"
    p2.font.size = Pt(11)
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER

def create_validation_slide(prs):
    """Create validation phase slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.6))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Phase 4: Validation - A/B Testing"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = VALIDATION_ORANGE

    # Table title
    table_title = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(9), Inches(0.4))
    tf = table_title.text_frame
    p = tf.paragraphs[0]
    p.text = "O'Brien-Fleming Early Stopping Boundaries"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = DARK_TEXT

    # Add table
    rows, cols = 5, 4
    table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.4), Inches(8), Inches(1.5)).table

    # Headers
    headers = ["Analysis", "Sample %", "Z-boundary", "p-boundary"]
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = VALIDATION_ORANGE
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = WHITE

    # Data
    data = [
        ["1st Interim", "25%", "4.05", "0.00005"],
        ["2nd Interim", "50%", "2.80", "0.0051"],
        ["3rd Interim", "75%", "2.28", "0.0226"],
        ["Final", "100%", "2.02", "0.0431"],
    ]

    for i, row_data in enumerate(data):
        for j, val in enumerate(row_data):
            cell = table.cell(i + 1, j)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(10)
            p.font.color.rgb = DARK_TEXT

    # Key points
    points_title = slide.shapes.add_textbox(Inches(0.5), Inches(3.1), Inches(9), Inches(0.4))
    tf = points_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Key Validation Components"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = DARK_TEXT

    points = [
        "✓ 2-week burn-in period for novelty effects",
        "✓ Segment analysis by psychological profiles",
        "✓ Guardrail metrics for safety monitoring",
        "✓ Bayesian alternative for flexible stopping"
    ]

    y = Inches(3.5)
    for point in points:
        box = slide.shapes.add_textbox(Inches(0.5), y, Inches(8), Inches(0.35))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = point
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_TEXT
        y += Inches(0.4)

def create_outcomes_slide(prs):
    """Create outcomes slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9), Inches(0.6))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Expected Outcomes"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = PRODUCTION_RED

    # Metrics boxes
    metrics = [
        ("-20%", "Cart Abandonment", "70% → 50%", PRODUCTION_RED),
        ("+10%", "Conversion Rate", "2.5% → 2.75%", PROTOTYPE_GREEN),
        ("-67%", "Checkout Time", "4.5min → 1.5min", RESEARCH_BLUE),
        ("+23%", "CSAT Score", "6.5 → 8.0", TRANSLATION_PURPLE),
    ]

    x_positions = [Inches(0.5), Inches(2.6), Inches(4.7), Inches(6.8)]

    for i, (change, label, detail, color) in enumerate(metrics):
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x_positions[i], Inches(1.2), Inches(2), Inches(2)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()

        tf = box.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = change
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = label
        p2.font.size = Pt(14)
        p2.font.bold = True
        p2.font.color.rgb = WHITE
        p2.alignment = PP_ALIGN.CENTER

        p3 = tf.add_paragraph()
        p3.text = detail
        p3.font.size = Pt(11)
        p3.font.color.rgb = WHITE
        p3.alignment = PP_ALIGN.CENTER

    # Summary
    summary = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(1.5))
    tf = summary.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Research-Driven Impact"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DARK_TEXT

    p2 = tf.add_paragraph()
    p2.text = "\nBy translating psychiatry research on cognitive load, decision-making, and reward systems into practical e-commerce features, we expect significant improvements across all key metrics."
    p2.font.size = Pt(14)
    p2.font.color.rgb = DARK_TEXT

def main():
    """Create the full presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9

    create_title_slide(prs)
    create_pipeline_slide(prs)
    create_research_phase_slide(prs)
    create_translation_phase_slide(prs)
    create_validation_slide(prs)
    create_outcomes_slide(prs)

    output_path = "docs/Research_To_Application_Pipeline.pptx"
    prs.save(output_path)
    print(f"✓ Presentation saved to: {output_path}")

if __name__ == "__main__":
    main()
