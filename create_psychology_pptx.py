#!/usr/bin/env python3
"""Create PowerPoint for Psychology → E-commerce Pipeline"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Colors
PSYCH_PURPLE = RGBColor(155, 89, 182)
BEHAVIOR_BLUE = RGBColor(52, 152, 219)
SOCIAL_GREEN = RGBColor(46, 204, 113)
COGNITIVE_ORANGE = RGBColor(230, 126, 34)
CONSUMER_RED = RGBColor(231, 76, 60)
DARK_TEXT = RGBColor(50, 50, 50)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(245, 245, 245)

def add_phase_box(slide, left, top, width, height, text, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    return shape

def add_subitem(slide, left, top, width, height, text, border_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = border_color
    shape.line.width = Pt(2)
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = DARK_TEXT
    p.alignment = PP_ALIGN.CENTER
    return shape

def create_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(1))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Psychology Research to E-Commerce"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = PSYCH_PURPLE
    p.alignment = PP_ALIGN.CENTER

    subtitle = slide.shapes.add_textbox(Inches(0.5), Inches(2.9), Inches(9), Inches(0.8))
    tf = subtitle.text_frame
    p = tf.paragraphs[0]
    p.text = "Translating Behavioral Science into Product Features"
    p.font.size = Pt(24)
    p.font.color.rgb = DARK_TEXT
    p.alignment = PP_ALIGN.CENTER

    sub2 = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(9), Inches(0.6))
    tf = sub2.text_frame
    p = tf.paragraphs[0]
    p.text = "Behavioral Commerce Platform"
    p.font.size = Pt(18)
    p.font.italic = True
    p.font.color.rgb = BEHAVIOR_BLUE
    p.alignment = PP_ALIGN.CENTER

    date = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(9), Inches(0.5))
    tf = date.text_frame
    p = tf.paragraphs[0]
    p.text = "January 2026 | Multi-Agent Research System"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(128, 128, 128)
    p.alignment = PP_ALIGN.CENTER

def create_pipeline_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title = slide.shapes.add_textbox(Inches(0.3), Inches(0.2), Inches(9.4), Inches(0.5))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Research to Application Pipeline"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = DARK_TEXT

    # Phase boxes
    phases = [
        ("1. PSYCHOLOGY\nRESEARCH", PSYCH_PURPLE, Inches(0.2)),
        ("2. INDUSTRY\nAPPLICATIONS", BEHAVIOR_BLUE, Inches(2.0)),
        ("3. PRODUCT\nSTRATEGY", SOCIAL_GREEN, Inches(3.8)),
        ("4. A/B\nTESTING", COGNITIVE_ORANGE, Inches(5.6)),
        ("5. PRODUCTION\nROLLOUT", CONSUMER_RED, Inches(7.4)),
    ]

    for text, color, left in phases:
        add_phase_box(slide, left, Inches(0.9), Inches(1.7), Inches(0.9), text, color)

    # Arrows
    for x in [Inches(1.95), Inches(3.75), Inches(5.55), Inches(7.35)]:
        shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, Inches(1.2), Inches(0.12), Inches(0.12))
        shape.fill.solid()
        shape.fill.fore_color.rgb = DARK_TEXT
        shape.line.fill.background()

    # Sub-items
    sub_y = Inches(2.0)
    sub_h = Inches(0.35)
    sub_w = Inches(1.6)
    gap = Inches(0.38)

    # Psychology Research items
    psych_items = ["Behavioral Psych", "Cognitive Psych", "Social Psych", "Consumer Psych", "UX Psychology"]
    for i, item in enumerate(psych_items):
        add_subitem(slide, Inches(0.25), sub_y + i*gap, sub_w, sub_h, item, PSYCH_PURPLE)

    # Industry items
    ind_items = ["Amazon", "Netflix", "Booking.com", "Social Proof", "Nudge Theory"]
    for i, item in enumerate(ind_items):
        add_subitem(slide, Inches(2.05), sub_y + i*gap, sub_w, sub_h, item, BEHAVIOR_BLUE)

    # Product items
    prod_items = ["User Personas", "Feature Specs", "MVP Definition", "Ethical Framework", "Go-to-Market"]
    for i, item in enumerate(prod_items):
        add_subitem(slide, Inches(3.85), sub_y + i*gap, sub_w, sub_h, item, SOCIAL_GREEN)

    # Testing items
    test_items = ["Sample Size", "Early Stopping", "Segmentation", "Metrics", "Duration"]
    for i, item in enumerate(test_items):
        add_subitem(slide, Inches(5.65), sub_y + i*gap, sub_w, sub_h, item, COGNITIVE_ORANGE)

    # Production items
    rollout_items = ["Phased Rollout", "Monitoring", "Iteration", "Scale", "Measurement"]
    for i, item in enumerate(rollout_items):
        add_subitem(slide, Inches(7.45), sub_y + i*gap, sub_w, sub_h, item, CONSUMER_RED)

    # Outcomes box
    outcomes = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.2), Inches(4.2), Inches(9.1), Inches(0.7))
    outcomes.fill.solid()
    outcomes.fill.fore_color.rgb = LIGHT_GRAY
    outcomes.line.color.rgb = DARK_TEXT
    tf = outcomes.text_frame
    p = tf.paragraphs[0]
    p.text = "Target: +20% Conversion | +40% Repeat Purchase | +33% Search Success | -14% Cart Abandonment"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = DARK_TEXT
    p.alignment = PP_ALIGN.CENTER

def create_psychology_domains_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(9), Inches(0.5))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Psychology Research Domains"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = PSYCH_PURPLE

    domains = [
        ("Behavioral Psychology", "Operant conditioning, habit formation, rewards", PSYCH_PURPLE),
        ("Cognitive Psychology", "Mental models, attention, memory, processing", BEHAVIOR_BLUE),
        ("Social Psychology", "Social proof, conformity, reciprocity", SOCIAL_GREEN),
        ("Consumer Psychology", "Price perception, anchoring, decision-making", COGNITIVE_ORANGE),
        ("UX Psychology", "Flow state, aesthetics, peak-end rule", CONSUMER_RED),
    ]

    y = Inches(0.75)
    for domain, desc, color in domains:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), y, Inches(9), Inches(0.75))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()

        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = domain
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = WHITE

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(12)
        p2.font.color.rgb = WHITE

        y += Inches(0.85)

def create_features_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(9), Inches(0.5))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Psychology-Based Features"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = SOCIAL_GREEN

    features = [
        ("Smart Reviews", "Social Proof", "Reviews, ratings, popularity indicators", "P0"),
        ("Quick Reorder", "Habit Formation", "One-click reorder, subscriptions, rewards", "P0"),
        ("Intelligent Search", "Cognitive Processing", "Autocomplete, suggestions, navigation", "P0"),
        ("Savings Display", "Price Anchoring", "Original price, % off, price alerts", "P1"),
        ("Urgency Indicators", "Scarcity Principle", "Stock levels, timers, FOMO triggers", "P1"),
    ]

    # Headers
    headers = ["Feature", "Psychology", "Components", "Priority"]
    x_positions = [Inches(0.3), Inches(2.1), Inches(4.0), Inches(8.2)]

    for i, (header, x) in enumerate(zip(headers, x_positions)):
        box = slide.shapes.add_textbox(x, Inches(0.7), Inches(1.8), Inches(0.35))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = header
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = DARK_TEXT

    y = Inches(1.1)
    for feature, psych, components, priority in features:
        # Feature name
        box1 = slide.shapes.add_textbox(Inches(0.3), y, Inches(1.7), Inches(0.6))
        tf = box1.text_frame
        p = tf.paragraphs[0]
        p.text = feature
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = DARK_TEXT

        # Psychology principle
        box2 = slide.shapes.add_textbox(Inches(2.1), y, Inches(1.8), Inches(0.6))
        tf = box2.text_frame
        p = tf.paragraphs[0]
        p.text = psych
        p.font.size = Pt(10)
        p.font.color.rgb = PSYCH_PURPLE

        # Components
        box3 = slide.shapes.add_textbox(Inches(4.0), y, Inches(4.1), Inches(0.6))
        tf = box3.text_frame
        p = tf.paragraphs[0]
        p.text = components
        p.font.size = Pt(10)
        p.font.color.rgb = DARK_TEXT

        # Priority
        prio_color = SOCIAL_GREEN if priority == "P0" else COGNITIVE_ORANGE
        box4 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.2), y, Inches(0.6), Inches(0.4))
        box4.fill.solid()
        box4.fill.fore_color.rgb = prio_color
        box4.line.fill.background()
        tf = box4.text_frame
        p = tf.paragraphs[0]
        p.text = priority
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        y += Inches(0.7)

def create_personas_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(9), Inches(0.5))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "User Personas by Psychology Profile"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = BEHAVIOR_BLUE

    personas = [
        ("Bargain Hunter Beth", "Price Anchoring", "Seeks deals, compares prices", PSYCH_PURPLE),
        ("Influence-Driven Ian", "Social Proof", "Reads reviews, follows trends", SOCIAL_GREEN),
        ("Habitual Helen", "Operant Conditioning", "Repeat buyer, values convenience", COGNITIVE_ORANGE),
        ("Occasional Owen", "Cognitive Efficiency", "Goal-oriented, needs fast search", CONSUMER_RED),
    ]

    positions = [
        (Inches(0.3), Inches(0.8)),
        (Inches(4.9), Inches(0.8)),
        (Inches(0.3), Inches(2.8)),
        (Inches(4.9), Inches(2.8)),
    ]

    for (name, psych, desc, color), (x, y) in zip(personas, positions):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(4.4), Inches(1.7))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()

        tf = box.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = WHITE

        p2 = tf.add_paragraph()
        p2.text = f"Psychology: {psych}"
        p2.font.size = Pt(12)
        p2.font.italic = True
        p2.font.color.rgb = WHITE

        p3 = tf.add_paragraph()
        p3.text = desc
        p3.font.size = Pt(11)
        p3.font.color.rgb = WHITE

def create_outcomes_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(9), Inches(0.5))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Expected Outcomes"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = CONSUMER_RED

    metrics = [
        ("+20%", "Conversion", "2.5% → 3.0%", PSYCH_PURPLE),
        ("+40%", "Repeat Purchase", "25% → 35%", SOCIAL_GREEN),
        ("+33%", "Search Success", "60% → 80%", BEHAVIOR_BLUE),
        ("-14%", "Cart Abandonment", "70% → 60%", CONSUMER_RED),
    ]

    x_positions = [Inches(0.4), Inches(2.55), Inches(4.7), Inches(6.85)]

    for (change, label, detail, color), x in zip(metrics, x_positions):
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1), Inches(2.1), Inches(2.2))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()

        tf = box.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = change
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = label
        p2.font.size = Pt(16)
        p2.font.bold = True
        p2.font.color.rgb = WHITE
        p2.alignment = PP_ALIGN.CENTER

        p3 = tf.add_paragraph()
        p3.text = detail
        p3.font.size = Pt(12)
        p3.font.color.rgb = WHITE
        p3.alignment = PP_ALIGN.CENTER

    # Summary
    summary = slide.shapes.add_textbox(Inches(0.4), Inches(3.5), Inches(9), Inches(1.2))
    tf = summary.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Psychology-Driven Impact"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DARK_TEXT

    p2 = tf.add_paragraph()
    p2.text = "By applying behavioral, cognitive, and social psychology principles to e-commerce design, we create more intuitive, trustworthy, and engaging shopping experiences."
    p2.font.size = Pt(13)
    p2.font.color.rgb = DARK_TEXT

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    create_title_slide(prs)
    create_pipeline_slide(prs)
    create_psychology_domains_slide(prs)
    create_personas_slide(prs)
    create_features_slide(prs)
    create_outcomes_slide(prs)

    output_path = "docs/Psychology_Ecommerce_Pipeline.pptx"
    prs.save(output_path)
    print(f"✓ Presentation saved to: {output_path}")

if __name__ == "__main__":
    main()
